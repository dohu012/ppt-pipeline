"""Generate a default template.pptx with all required slide layouts.

This is a convenience utility so the pipeline works out-of-the-box without
requiring the user to manually craft a template in PowerPoint.
Users are encouraged to replace the generated template with their own design.
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _add_placeholder(slide_layout, idx: int, name: str, left, top, width, height):
    """Add a placeholder shape to a slide layout."""
    # python-pptx doesn't expose a clean API for adding placeholders to
    # slide layouts, so we work at the XML level.
    from lxml import etree

    NSMAP = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    sp_tree = slide_layout.element

    # Find or create the cSld/spTree element
    cSld = sp_tree.find(".//p:cSld", NSMAP)
    if cSld is None:
        cSld = etree.SubElement(
            sp_tree,
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld",
        )
    spTree = cSld.find(".//p:spTree", NSMAP)
    if spTree is None:
        spTree = etree.SubElement(
            cSld,
            "{http://schemas.openxmlformats.org/presentationml/2006/main}spTree",
        )

    # Build a <p:sp> shape element
    shape_elm = etree.SubElement(spTree, "{http://schemas.openxmlformats.org/presentationml/2006/main}sp")
    shape_elm.set("useBgFill", "1")

    nvSpPr = etree.SubElement(shape_elm, "{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr")
    cNvPr = etree.SubElement(nvSpPr, "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
    cNvPr.set("id", str(idx))
    cNvPr.set("name", name)
    cNvSpPr = etree.SubElement(nvSpPr, "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvSpPr")
    cNvSpPr.set("txBox", "1")
    nvPr = etree.SubElement(nvSpPr, "{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr")

    # <p:spPr>
    spPr = etree.SubElement(shape_elm, "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
    xfrm = etree.SubElement(spPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
    off = etree.SubElement(xfrm, "{http://schemas.openxmlformats.org/drawingml/2006/main}off")
    off.set("x", str(int(left)))
    off.set("y", str(int(top)))
    ext = etree.SubElement(xfrm, "{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
    ext.set("cx", str(int(width)))
    ext.set("cy", str(int(height)))
    prstGeom = etree.SubElement(spPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
    prstGeom.set("prst", "rect")
    avLst = etree.SubElement(prstGeom, "{http://schemas.openxmlformats.org/drawingml/2006/main}avLst")

    # <p:txBody>
    txBody = etree.SubElement(shape_elm, "{http://schemas.openxmlformats.org/presentationml/2006/main}txBody")
    bodyPr = etree.SubElement(txBody, "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
    p_elm = etree.SubElement(txBody, "{http://schemas.openxmlformats.org/drawingml/2006/main}p")

    return shape_elm


def create_default_template(output_path: str = "template.pptx") -> str:
    """Generate a basic template.pptx with all 7 layout types.

    Layout names follow the ppt_plan.json spec:
    title, toc, section_title, bullets, figure, table, end
    """
    prs = Presentation()

    # We want 16:9 slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def make_layout(name: str) -> None:
        layout = prs.slide_layouts.add_new_layout()
        layout.name = name
        # The built-in blank layout has minimal decorations
        return layout

    # ------------------------------------------------------------------
    # title — 封面
    # ------------------------------------------------------------------
    title_layout = prs.slide_layouts[0]  # reuse the default title layout
    title_layout.name = "title"
    # Its default placeholders: title (idx 0), subtitle (idx 1)
    # Rename placeholders for clarity
    for ph in title_layout.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.name = "title"
        elif ph.placeholder_format.idx == 1:
            ph.name = "subtitle"

    # ------------------------------------------------------------------
    # toc — 目录
    # ------------------------------------------------------------------
    toc_layout = prs.slide_layouts.add_new_layout()
    toc_layout.name = "toc"
    _add_placeholder(toc_layout, 0, "title", Inches(1), Inches(0.5), Inches(11), Inches(1))
    _add_placeholder(toc_layout, 1, "content", Inches(2), Inches(2), Inches(9), Inches(4.5))

    # ------------------------------------------------------------------
    # section_title — 章节过渡页
    # ------------------------------------------------------------------
    sec_layout = prs.slide_layouts.add_new_layout()
    sec_layout.name = "section_title"
    _add_placeholder(sec_layout, 0, "title", Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    _add_placeholder(sec_layout, 1, "subtitle", Inches(1.5), Inches(4.2), Inches(10), Inches(1))

    # ------------------------------------------------------------------
    # bullets — 正文内容
    # ------------------------------------------------------------------
    bullets_layout = prs.slide_layouts.add_new_layout()
    bullets_layout.name = "bullets"
    _add_placeholder(bullets_layout, 0, "title", Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    _add_placeholder(bullets_layout, 1, "body", Inches(1.2), Inches(1.6), Inches(10), Inches(5.2))

    # ------------------------------------------------------------------
    # figure — 图表展示
    # ------------------------------------------------------------------
    figure_layout = prs.slide_layouts.add_new_layout()
    figure_layout.name = "figure"
    _add_placeholder(figure_layout, 0, "title", Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    _add_placeholder(figure_layout, 1, "image", Inches(1), Inches(1.5), Inches(11), Inches(5))
    _add_placeholder(figure_layout, 2, "caption", Inches(1), Inches(6.7), Inches(11), Inches(0.6))

    # ------------------------------------------------------------------
    # table — 数据对比
    # ------------------------------------------------------------------
    table_layout = prs.slide_layouts.add_new_layout()
    table_layout.name = "table"
    _add_placeholder(table_layout, 0, "title", Inches(0.8), Inches(0.4), Inches(11), Inches(0.9))
    _add_placeholder(table_layout, 1, "table", Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))

    # ------------------------------------------------------------------
    # end — 结尾页
    # ------------------------------------------------------------------
    end_layout = prs.slide_layouts.add_new_layout()
    end_layout.name = "end"
    _add_placeholder(end_layout, 0, "text", Inches(2), Inches(3), Inches(9), Inches(1.5))

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    out = sys.argv[1] if len(sys.argv) > 1 else "template.pptx"
    create_default_template(out)
    print(f"Default template created → {out}")
