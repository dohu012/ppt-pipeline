"""Phase 2: Render ppt_plan.json into a .pptx file using a template."""

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ---------------------------------------------------------------------------
# Layout name → slide layout lookup
# ---------------------------------------------------------------------------

_LAYOUT_CACHE: dict[str, int] = {}

# Fallback placeholder names for each layout type when template doesn't
# use standard PowerPoint placeholder names.
_PLACEHOLDER_FALLBACKS: dict[str, dict[str, str]] = {
    "title": {"title": "标题", "subtitle": "副标题", "author": "作者", "date": "日期"},
    "toc": {"title": "标题", "items": "内容"},
    "section_title": {"title": "标题", "subtitle": "副标题"},
    "bullets": {"title": "标题", "body": "正文"},
    "figure": {"title": "标题", "image": "图片", "caption": "图注"},
    "table": {"title": "标题", "table": "表格"},
    "end": {"text": "致谢"},
}

# English → Chinese synonyms for placeholder name matching
_KEYWORD_SYNONYMS: dict[str, list[str]] = {
    "title":    ["title", "标题"],
    "subtitle": ["subtitle", "副标题"],
    "author":   ["author", "作者"],
    "date":     ["date", "日期"],
    "body":     ["body", "正文", "内容"],
    "items":    ["items", "内容", "列表", "项目"],
    "content":  ["content", "内容", "正文"],
    "image":    ["image", "图片", "picture", "img"],
    "caption":  ["caption", "图注", "subtitle"],
    "table":    ["table", "表格"],
    "text":     ["text", "文本", "文字", "正文", "内容"],
    "thanks":   ["thanks", "致谢"],
}

def _build_layout_map(prs: Presentation) -> dict[str, int]:
    """Match each slide type to the best layout using name + shape analysis."""
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        n_text = sum(1 for s in layout.shapes if s.has_text_frame)
        name_lower = layout.name.lower()
        layouts.append({
            "idx": i, "name": layout.name, "n_text": n_text,
            "is_cover":  any(kw in name_lower for kw in ["封面", "cover"]),
            "is_end":    any(kw in name_lower for kw in ["尾页", "封底", "end"]),
            "is_section": any(kw in name_lower for kw in ["章节", "section"]),
            "is_blank":  any(kw in name_lower for kw in ["空白", "blank"]),
        })

    mapping = {}

    # title → cover layout
    cover = [ly for ly in layouts if ly["is_cover"]]
    mapping["title"] = cover[0]["idx"] if cover else 0

    # end → end/back-cover layout
    end = [ly for ly in layouts if ly["is_end"]]
    mapping["end"] = end[0]["idx"] if end else 0

    # section_title → section layout
    sec = [ly for ly in layouts if ly["is_section"]]
    mapping["section_title"] = sec[0]["idx"] if sec else 0

    # bullets/toc/figure/table → blank (or first with most text shapes)
    blank = [ly for ly in layouts if ly["is_blank"]]
    if blank:
        fallback = blank[0]["idx"]
    else:
        layouts.sort(key=lambda x: x["n_text"], reverse=True)
        fallback = layouts[0]["idx"]
    for st in ("bullets", "toc", "figure", "table"):
        mapping[st] = fallback

    print(f"  Layout mapping: { {k: layouts[v]['name'] for k, v in mapping.items()} }")
    return mapping


def _find_layout_index(prs: Presentation, name: str) -> int:
    """Find a slide layout by name, with dynamic fallback."""
    cache_key = f"{id(prs)}_{name}"
    if cache_key in _LAYOUT_CACHE:
        return _LAYOUT_CACHE[cache_key]

    # Dynamic matching: build once per presentation
    dynamic_key = f"{id(prs)}_dynamic"
    if dynamic_key not in _LAYOUT_CACHE:
        layout_map = _build_layout_map(prs)
        for k, v in layout_map.items():
            _LAYOUT_CACHE[f"{id(prs)}_{k}"] = v
        _LAYOUT_CACHE[dynamic_key] = True

    if f"{id(prs)}_{name}" in _LAYOUT_CACHE:
        return _LAYOUT_CACHE[f"{id(prs)}_{name}"]

    # Final fallback: first layout
    _LAYOUT_CACHE[cache_key] = 0
    return 0


# ---------------------------------------------------------------------------
# Placeholder helpers
# ---------------------------------------------------------------------------

def _set_text(shape, text: str) -> None:
    """Write text into a shape, handling both auto-shapes and text frames."""
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
    elif shape.has_table:
        shape.table.cell(0, 0).text = text


def _find_placeholder(slide, keywords: list[str]):
    """Find a shape by name keyword matching (with Chinese/English synonyms)."""
    expanded = []
    for kw in keywords:
        expanded.append(kw)
        expanded.extend(_KEYWORD_SYNONYMS.get(kw, []))
    seen: set[str] = set()
    unique = []
    for kw in expanded:
        kwl = kw.lower()
        if kwl not in seen:
            seen.add(kwl)
            unique.append(kwl)

    for shape in slide.shapes:
        name_lower = shape.name.lower()
        for kw in unique:
            if kw in name_lower:
                return shape
    return None


def _find_body_shape(slide):
    """Find the best shape for body text on a slide.

    1. Named body/content placeholder → use it
    2. Largest text-bearing shape (excluding title) → use it
    3. Fallback → create a text box
    """
    body = _find_placeholder(slide, ["body", "content", "text", "bullets", "items"])
    if body:
        return body

    title = _find_placeholder(slide, ["title"])
    candidates = []
    for shape in slide.shapes:
        if shape == title or not shape.has_text_frame:
            continue
        candidates.append((shape.width * shape.height, shape))

    if candidates:
        candidates.sort(key=lambda x: x[0], reverse=True)
        return candidates[0][1]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    return txBox


def _add_bullet_paragraph(tf, text: str, level: int = 0, font_size: int = 14) -> None:
    """Append a paragraph with bullet formatting to a text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.space_after = Pt(4)
    for run in p.runs:
        run.font.size = Pt(font_size)


# ---------------------------------------------------------------------------
# Slide fillers
# ---------------------------------------------------------------------------

def _fill_title(slide, content: dict) -> None:
    for key in ("title", "subtitle", "author", "date"):
        val = content.get(key, "")
        if not val:
            continue
        ph = _find_placeholder(slide, [key])
        if ph:
            _set_text(ph, val)


def _fill_toc(slide, content: dict) -> None:
    # Title
    ph = _find_placeholder(slide, ["title"])
    if ph:
        _set_text(ph, content.get("title", "目录"))

    # Items list
    items = content.get("items", [])
    body = _find_body_shape(slide)
    if body and body.has_text_frame:
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = items[0] if items else ""
        for item in items[1:]:
            _add_bullet_paragraph(tf, item)


def _fill_section_title(slide, content: dict) -> None:
    for key in ("title", "subtitle"):
        val = content.get(key, "")
        if not val:
            continue
        ph = _find_placeholder(slide, [key])
        if ph:
            _set_text(ph, val)


def _fill_bullets(slide, content: dict) -> None:
    # Title
    ph = _find_placeholder(slide, ["title"])
    if ph:
        _set_text(ph, content.get("title", ""))

    # Bullet list
    bullets = content.get("bullets", [])
    body = _find_body_shape(slide)
    if body and body.has_text_frame:
        tf = body.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            text = b["bullet"] if isinstance(b, dict) else str(b)
            if i == 0:
                tf.paragraphs[0].text = text
            else:
                _add_bullet_paragraph(tf, text)


def _fill_figure(slide, content: dict) -> None:
    # Title
    ph = _find_placeholder(slide, ["title"])
    if ph:
        _set_text(ph, content.get("title", ""))

    # Image
    img_path = content.get("image", "")
    if img_path and Path(img_path).exists():
        img_ph = _find_placeholder(slide, ["image", "picture", "图片", "img"])
        if img_ph:
            # Replace placeholder with picture
            left = img_ph.left
            top = img_ph.top
            width = img_ph.width
            height = img_ph.height
            sp = img_ph._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(str(img_path), left, top, width, height)
        else:
            # Just add picture centered on slide
            slide.shapes.add_picture(
                str(img_path),
                Inches(1.5),
                Inches(2),
                Inches(7),
                Inches(4.5),
            )

    # Caption
    caption = content.get("caption", "")
    if caption:
        ph = _find_placeholder(slide, ["caption", "图注", "subtitle"])
        if ph:
            _set_text(ph, caption)


def _fill_table(slide, content: dict) -> None:
    # Title
    ph = _find_placeholder(slide, ["title"])
    if ph:
        _set_text(ph, content.get("title", ""))

    header = content.get("header", [])
    rows = content.get("rows", [])

    if not header or not rows:
        return

    tbl_ph = _find_placeholder(slide, ["table", "表格"])
    if tbl_ph:
        left = tbl_ph.left
        top = tbl_ph.top
        width = tbl_ph.width
        height = tbl_ph.height
        sp = tbl_ph._element
        sp.getparent().remove(sp)
    else:
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(8.4)
        height = Inches(4.5)

    n_rows = len(rows) + 1
    n_cols = len(header)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header row
    for ci, col_name in enumerate(header):
        cell = table.cell(0, ci)
        cell.text = col_name
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.bold = True

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            if ci < n_cols:
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(11)


def _fill_end(slide, content: dict) -> None:
    text = content.get("text", "感谢各位老师批评指正")
    ph = _find_placeholder(slide, ["text", "致谢", "thanks", "content", "body", "title"])
    if ph:
        _set_text(ph, text)


# ---------------------------------------------------------------------------
# Filler dispatch
# ---------------------------------------------------------------------------

_FILLERS = {
    "title": _fill_title,
    "toc": _fill_toc,
    "section_title": _fill_section_title,
    "bullets": _fill_bullets,
    "figure": _fill_figure,
    "table": _fill_table,
    "end": _fill_end,
}


# ---------------------------------------------------------------------------
# Main renderer
# ---------------------------------------------------------------------------

def render(
    plan_path: str | Path,
    template_path: str | Path = "template.pptx",
    output_path: str | Path = "output/答辩PPT.pptx",
) -> str:
    """Render ppt_plan.json into a .pptx file.

    Args:
        plan_path: Path to ppt_plan.json.
        template_path: Path to the .pptx template with predefined layouts.
        output_path: Where to write the final .pptx.

    Returns:
        The output file path as a string.
    """
    plan_path = Path(plan_path)
    template_path = Path(template_path)
    output_path = Path(output_path)

    if not plan_path.exists():
        raise FileNotFoundError(f"Plan file not found: {plan_path}")
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    prs = Presentation(str(template_path))

    for slide_plan in plan["slides"]:
        layout_name = slide_plan["layout"]
        layout_idx = _find_layout_index(prs, layout_name)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        filler = _FILLERS.get(layout_name)
        if filler:
            filler(slide, slide_plan["content"])
        else:
            print(f"  [warn] Unknown layout '{layout_name}' — skipping content fill")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    plan_file = sys.argv[1] if len(sys.argv) > 1 else "output/ppt_plan.json"
    tmpl_file = sys.argv[2] if len(sys.argv) > 2 else "template.pptx"
    out_file = sys.argv[3] if len(sys.argv) > 3 else "output/答辩PPT.pptx"

    result = render(plan_file, tmpl_file, out_file)
    print(f"PPTX rendered → {result}")
